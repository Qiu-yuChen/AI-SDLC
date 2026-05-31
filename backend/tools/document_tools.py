"""Document Parsing Tools — Extract text content from multiple file formats

Provides pure functions (not CrewAI tools) for parsing various document formats
into structured plain text, ready for LLM-based intent recognition.

Supported formats: md, docx, pptx, pdf, xlsx, csv, txt, html, json, xml,
                    yaml/yml, rst, org, adoc
"""

import io
import csv as csv_module
from pathlib import Path
from typing import Optional, Dict, List


# ═══════════════════════════════════════════════════════════════
#  Individual Format Parsers
# ═══════════════════════════════════════════════════════════════

def parse_docx(file_path: str) -> Dict[str, str]:
    """Parse .docx into structured text with headings, body, and tables."""
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx", "raw": "", "title": ""}

    try:
        doc = Document(file_path)
    except Exception as e:
        return {"error": f"Failed to open .docx: {e}", "raw": "", "title": ""}

    paragraphs = []
    headings = []
    tables_text = []
    title = ""

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = para.style.name.lower() if para.style else ""

        if "heading 1" in style_name or "title" in style_name:
            if not title:
                title = text
            headings.append(f"# {text}")
            paragraphs.append(f"\n# {text}\n")
        elif "heading 2" in style_name:
            headings.append(f"  ## {text}")
            paragraphs.append(f"\n## {text}\n")
        elif "heading 3" in style_name:
            headings.append(f"    ### {text}")
            paragraphs.append(f"\n### {text}\n")
        elif "heading" in style_name:
            headings.append(f"      #### {text}")
            paragraphs.append(f"\n#### {text}\n")
        else:
            paragraphs.append(text)

    for table_idx, table in enumerate(doc.tables):
        table_lines = [f"\n### 表格 {table_idx + 1}\n"]
        header = [cell.text.strip() for cell in table.rows[0].cells]
        table_lines.append("| " + " | ".join(header) + " |")
        table_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            table_lines.append("| " + " | ".join(cells) + " |")
        table_md = "\n".join(table_lines)
        tables_text.append(table_md)
        paragraphs.append(table_md)

    body = "\n\n".join(paragraphs)
    return {
        "title": title or Path(file_path).stem,
        "headings": "\n".join(headings) if headings else "(无标题结构)",
        "body": body,
        "tables": "\n\n".join(tables_text) if tables_text else "(无表格)",
        "raw": body,
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
    }


def parse_pptx(file_path: str) -> Dict[str, str]:
    """Parse .pptx into structured text with slides and speaker notes."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx", "raw": "", "title": ""}

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"error": f"Failed to open .pptx: {e}", "raw": "", "title": ""}

    slides_text = []
    title = ""

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_body = []
        slide_notes = ""

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.type == 1:
                    slide_title = shape.text.strip()
                elif ph.type in (2, 3, 7):
                    slide_body.append(shape.text.strip())
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if not slide_title and slide_num == 1:
                        slide_title = text
                    else:
                        slide_body.append(text)
            if shape.has_table:
                table = shape.table
                table_lines = ["\n| " + " | ".join(
                    cell.text.strip() for cell in row.cells
                ) + " |" for row in table.rows]
                slide_body.append("\n".join(table_lines))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_notes = slide.notes_slide.notes_text_frame.text.strip()

        if slide_num == 1 and slide_title:
            title = slide_title

        slide_md = f"## Slide {slide_num}: {slide_title or '(无标题)'}\n"
        if slide_body:
            slide_md += "\n".join(slide_body) + "\n"
        if slide_notes:
            slide_md += f"\n> **备注**: {slide_notes}\n"
        slides_text.append(slide_md)

    raw = "\n\n---\n\n".join(slides_text)
    return {
        "title": title or Path(file_path).stem,
        "slides": slides_text,
        "raw": raw,
        "slide_count": len(prs.slides),
    }


def parse_pdf(file_path: str) -> Dict[str, str]:
    """Parse .pdf into plain text, extracting all pages."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return {"error": "PyPDF2 not installed. Run: pip install PyPDF2", "raw": "", "title": ""}

    try:
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"## 第 {i + 1} 页\n\n{text.strip()}")
        raw = "\n\n---\n\n".join(pages)
        title = Path(file_path).stem
        return {
            "title": title,
            "raw": raw,
            "page_count": len(reader.pages),
        }
    except Exception as e:
        return {"error": f"Failed to parse PDF: {e}", "raw": "", "title": Path(file_path).stem}


def parse_xlsx(file_path: str) -> Dict[str, str]:
    """Parse .xlsx into Markdown tables, one per sheet."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"error": "openpyxl not installed. Run: pip install openpyxl", "raw": "", "title": ""}

    try:
        wb = load_workbook(file_path, data_only=True)
        sheets_md = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            lines = [f"## Sheet: {sheet_name}\n"]
            # Build markdown table
            max_cols = max(len(row) for row in rows) if rows else 0
            header = [str(c) if c is not None else "" for c in rows[0]]
            # Pad header if needed
            while len(header) < max_cols:
                header.append("")
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * max_cols) + " |")
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]
                while len(cells) < max_cols:
                    cells.append("")
                lines.append("| " + " | ".join(cells) + " |")
            sheets_md.append("\n".join(lines))

        raw = "\n\n".join(sheets_md)
        return {
            "title": Path(file_path).stem,
            "raw": raw,
            "sheet_count": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
        }
    except Exception as e:
        return {"error": f"Failed to parse xlsx: {e}", "raw": "", "title": Path(file_path).stem}


def parse_csv(file_path: str) -> Dict[str, str]:
    """Parse .csv into Markdown table."""
    try:
        with open(file_path, "r", encoding="utf-8-sig") as f:
            reader = csv_module.reader(f)
            rows = list(reader)
        if not rows:
            return {"title": Path(file_path).stem, "raw": "(空文件)", "row_count": 0}

        max_cols = max(len(row) for row in rows)
        lines = []
        header = rows[0]
        while len(header) < max_cols:
            header.append("")
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in rows[1:]:
            cells = list(row)
            while len(cells) < max_cols:
                cells.append("")
            lines.append("| " + " | ".join(cells) + " |")

        raw = "\n".join(lines)
        return {
            "title": Path(file_path).stem,
            "raw": raw,
            "row_count": len(rows),
        }
    except Exception as e:
        return {"error": f"Failed to parse CSV: {e}", "raw": "", "title": Path(file_path).stem}


def parse_html(file_path: str) -> Dict[str, str]:
    """Parse .html into plain text (strips tags)."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Fallback: basic tag stripping
        return _parse_html_fallback(file_path)

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        # Remove script/style tags
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        title = soup.title.string.strip() if soup.title else Path(file_path).stem
        raw = soup.get_text("\n", strip=True)
        return {"title": title, "raw": raw}
    except Exception as e:
        return _parse_html_fallback(file_path)


def _parse_html_fallback(file_path: str) -> Dict[str, str]:
    """Basic HTML-to-text without BeautifulSoup."""
    import re
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        # Strip tags
        text = re.sub(r"<[^>]+>", " ", content)
        text = re.sub(r"\s+", " ", text).strip()
        # Try to extract title
        title_match = re.search(r"<title[^>]*>(.*?)</title>", content, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else Path(file_path).stem
        return {"title": title, "raw": text}
    except Exception as e:
        return {"error": f"Failed to parse HTML: {e}", "raw": "", "title": Path(file_path).stem}


def parse_json(file_path: str) -> Dict[str, str]:
    """Parse .json into pretty-printed structure overview."""
    import json as json_module
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        data = json_module.loads(content)

        # Build a readable summary
        lines = [f"```json\n{json_module.dumps(data, ensure_ascii=False, indent=2)}\n```"]

        # Add structure summary if it's a dict
        if isinstance(data, dict):
            lines.insert(0, f"## JSON 顶层键: {', '.join(data.keys())}")
            lines.insert(0, f"# {Path(file_path).stem}")

        raw = "\n\n".join(lines)
        return {"title": Path(file_path).stem, "raw": raw}
    except Exception as e:
        return {"error": f"Failed to parse JSON: {e}", "raw": "", "title": Path(file_path).stem}


def parse_xml(file_path: str) -> Dict[str, str]:
    """Parse .xml into indented text representation."""
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        # Try to pretty-print
        try:
            import xml.dom.minidom as minidom
            dom = minidom.parseString(content)
            pretty = dom.toprettyxml(indent="  ")
        except Exception:
            pretty = content
        return {"title": Path(file_path).stem, "raw": f"```xml\n{pretty}\n```"}
    except Exception as e:
        return {"error": f"Failed to parse XML: {e}", "raw": "", "title": Path(file_path).stem}


def parse_yaml(file_path: str) -> Dict[str, str]:
    """Parse .yaml/.yml into structured text."""
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        # Try to load and re-dump for structure overview
        try:
            import yaml
            data = yaml.safe_load(content)
            import json
            pretty = json.dumps(data, ensure_ascii=False, indent=2)
            return {"title": Path(file_path).stem, "raw": f"```json\n{pretty}\n```"}
        except ImportError:
            pass
        return {"title": Path(file_path).stem, "raw": f"```yaml\n{content}\n```"}
    except Exception as e:
        return {"error": f"Failed to parse YAML: {e}", "raw": "", "title": Path(file_path).stem}


def parse_rst(file_path: str) -> Dict[str, str]:
    """Parse .rst (reStructuredText) into plain text."""
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        # Try to extract title from first header
        title = Path(file_path).stem
        for line in content.split("\n"):
            stripped = line.strip()
            if stripped and all(c in "=-~^\"'`*+#" for c in stripped):
                continue  # underline/overline decoration
            if stripped:
                title = stripped
                break
        # Basic RST-to-text: just return as-is (LLM can handle rst markup)
        return {"title": title, "raw": content}
    except Exception as e:
        return {"error": f"Failed to read RST: {e}", "raw": "", "title": Path(file_path).stem}


# ── Simple text-based formats (direct read) ──────────────────

def _parse_text_file(file_path: str) -> Dict[str, str]:
    """Generic parser for plain-text formats: .txt, .org, .adoc, .log, etc."""
    try:
        content = Path(file_path).read_text(encoding="utf-8")
        # Try to extract first non-empty meaningful line as title
        title = Path(file_path).stem
        for line in content.split("\n"):
            stripped = line.strip()
            if stripped and not stripped.startswith("```"):
                # For .org files, look for #+TITLE:
                if stripped.lower().startswith("#+title:"):
                    title = stripped.split(":", 1)[1].strip()
                    break
                # For .adoc files, look for = Title
                if stripped.startswith("= ") and not stripped.startswith("== "):
                    title = stripped[2:].strip()
                    break
                title = stripped[:80]
                break
        return {"title": title, "raw": content}
    except Exception as e:
        return {"error": f"Failed to read file: {e}", "raw": "", "title": Path(file_path).stem}


# ═══════════════════════════════════════════════════════════════
#  Extension Registry & Dispatcher
# ═══════════════════════════════════════════════════════════════

SUPPORTED_EXTENSIONS = {
    ".md": "markdown",
    ".docx": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".csv": "csv",
    ".txt": "text",
    ".html": "html",
    ".htm": "html",
    ".json": "json",
    ".xml": "xml",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".rst": "rst",
    ".org": "org",
    ".adoc": "adoc",
    ".log": "text",
}

# Text-based formats that don't need a special parser
_TEXT_FORMATS = {".txt", ".org", ".adoc", ".log", ".rst"}


def get_file_type(filename: str) -> Optional[str]:
    """Determine document type from filename extension."""
    ext = Path(filename).suffix.lower()
    return SUPPORTED_EXTENSIONS.get(ext)


def parse_document(file_path: str) -> Dict[str, str]:
    """
    Auto-detect file type and parse accordingly.
    Returns a dict with at least 'raw' and 'title' keys.
    """
    ext = Path(file_path).suffix.lower()

    # ── Structured format parsers ──
    if ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext in (".html", ".htm"):
        return parse_html(file_path)
    elif ext == ".json":
        return parse_json(file_path)
    elif ext == ".xml":
        return parse_xml(file_path)
    elif ext in (".yaml", ".yml"):
        return parse_yaml(file_path)

    # ── Text-based formats (direct read) ──
    elif ext in _TEXT_FORMATS:
        return _parse_text_file(file_path)

    # ── Fallback: read as plain text ──
    else:
        return _parse_text_file(file_path)
